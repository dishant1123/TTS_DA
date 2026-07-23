# regression  plot  , adv style   in seaborn:

import seaborn as sns
import matplotlib.pyplot as plt
import pandas as pd


df=pd.read_csv("seaborn/sales_discount.csv")
print(df.head())

# plt.figure(figsize=(10,5))

"""sns.regplot(
    data=df,
    x="Discount",
    y="Sales",
    ci=None,  # confidence intervals,
    scatter_kws={"s": 100}
)
plt.title("Sales vs Discount")
plt.show()

"""

# lmplot :
"""
sns.lmplot(
    data=df,
    x="Discount",
    y="Sales",
    col="Region",
    height=4

)
plt.show()
"""

# adv style :1. sns.set_theme() sns.set_palette(),sns.set_context()

# sns.set_theme(
# sns.set_theme(style="darkgrid")
# sns.set_theme(style="whitegrid")
"""sns.set_theme(style="ticks")



sns.regplot(
    data=df,
    x="Discount",
    y="Sales",
    ci=None,  # confidence intervals,
)
plt.show()

"""

# color palette :

"""# sns.set_theme(style="ticks")
# sns.set_palette("deep")
# sns.set_palette("pastel", n_colors=5)
sns.set_palette("colorblind")


sns.regplot(
    data=df,
    x="Discount",
    y="Sales"
    # ci=None  # confidence intervals,
)
plt.show()

"""

# sns.set_context : 

"""
# sns.set_context("paper")
# sns.set_context("talk")
# sns.set_context("poster")
sns.set_context("notebook")

sns.regplot(
    data=df,
    x="Discount",
    y="Sales"
    # ci=None  # confidence intervals,
)
plt.show()
"""

# corporate style :

"""sns.set_theme(
    style="whitegrid",
    palette="deep",
    context="talk"
)

plt.figure(figsize=(10,5))

sns.regplot(
    data=df,
    x="Discount",
    y="Sales",
    scatter_kws={"s": 100},
    line_kws={"linewidth": 3}
)
plt.title("Sales vs Discount")
plt.xlabel("Discount")
plt.ylabel("Sales")
plt.show()
"""

# pdf ,svg , png , jpg  : 

plt.figure(figsize=(10,5))

sns.scatterplot(
    data=df,
    x="Discount",
    y="Sales",
    hue="Region",
    s=120
)
plt.title("Sales vs Discount")

plt.savefig("seaborn/sales_discount.png")  # after  save  then  plt.show() use  

# plt.savefig("seaborn/SALES_DISCOUNT_HIGH_QUALITY.png",dpi=300)   # dpi  use  for  picture quality 

# plt.savefig("seaborn/sales_discount.pdf") 

# plt.savefig("seaborn/sales_discount_transparent.png",transparent=True)   # transparent  use  for  picture quality

# plt.savefig("seaborn/salesdis_analysis.png",bbox_inches="tight")   # bbox_inches  use  for  picture quality


plt.show()

# ppt : 
"""
pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()


slide_layout=prs.slide_layouts[5]
slide =prs.slides.add_slide(slide_layout)

slide.shapes.title.text = "Sales vs Discount"

slide.shapes.add_picture(
    "seaborn/sales_discount.png",
    Inches(1.5), 
    Inches(1.5),
    width=Inches(7),
)

prs.save("sales_discount.pptx")
print("pptx file saved successfully")


# 1 aug  submit all  assignment for numpy pandas matplotlib seaborn 